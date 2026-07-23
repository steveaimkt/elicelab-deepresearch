#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
조사결과.md  →  WMBB deck 디자인 PPTX (4교시 최종 산출물)
------------------------------------------------------------------------
핵심: PPTX는 조사결과.md에서 자동으로 정리된다. md를 고치면 PPTX도 따라간다.

쓰는 법:
  python3 보고서_pptx_템플릿.py [조사결과.md 경로]
  (경로 생략 시 같은 폴더의 조사결과.md 사용)  →  조사결과.pptx 생성.
필요: python-pptx

조사결과.md 구조(에이전트가 이 형식으로 저장):
  # {제목}
  메타: ...
  질문: ...
  결론: ...
  결론보충: ...
  ## 인사이트
  - [반전] {한 줄} :: {보충}
  ## 반대 근거
  반대핵심: ... / 반대문장: ... / 반대보충: ...
  ## 액션
  - {액션}
  ## 출처
  출처요약: ... / - {출처} / 버림: ...
"""
import sys, os, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

# ─────────────────────── 조사결과.md 파싱 ───────────────────────
def parse_md(path):
    lines=open(path,encoding="utf-8").read().splitlines()
    d={"title":"리서치 보고서","meta":"","question":"","conclusion":"","conclusion_sub":"",
       "insights":[],"counter_hl":"","counter":"","counter_sub":"",
       "actions":[],"sources_head":"출처 목록","sources":[],"dropped":""}
    sec=None
    def val(l): return l.split(":",1)[1].strip()
    for l in lines:
        t=l.strip()
        if t.startswith("# "):        d["title"]=t[2:].strip(); continue
        if t.startswith("## "):       sec=t[3:].strip(); continue
        if t.startswith("메타:"):      d["meta"]=val(t)
        elif t.startswith("질문:"):    d["question"]=val(t)
        elif t.startswith("결론보충:"): d["conclusion_sub"]=val(t)
        elif t.startswith("결론:"):    d["conclusion"]=val(t)
        elif t.startswith("반대핵심:"): d["counter_hl"]=val(t)
        elif t.startswith("반대문장:"): d["counter"]=val(t)
        elif t.startswith("반대보충:"): d["counter_sub"]=val(t)
        elif t.startswith("출처요약:"): d["sources_head"]=val(t)
        elif t.startswith("버림:"):    d["dropped"]=t
        elif t.startswith("- "):
            item=t[2:].strip()
            if sec=="인사이트":
                m=re.match(r"\[(.+?)\]\s*(.+?)\s*::\s*(.+)",item)
                if m: d["insights"].append((m.group(1),m.group(2),m.group(3)))
                else: d["insights"].append(("", item, ""))
            elif sec=="액션": d["actions"].append(item)
            elif sec=="출처": d["sources"].append(item)
    return d

md_path = sys.argv[1] if len(sys.argv)>1 else os.path.join(os.path.dirname(os.path.abspath(__file__)),"조사결과.md")
if not os.path.exists(md_path):
    md_path=os.path.join(os.getcwd(),"조사결과.md")
C=parse_md(md_path)
OUT=os.path.join(os.path.dirname(md_path),"조사결과.pptx")

# ─────────────────────── WMBB deck 렌더 ───────────────────────
INK=RGBColor(0x2A,0x2A,0x28); PAPER=RGBColor(0xF3,0xF3,0xF2)
WHITE=RGBColor(0xFF,0xFF,0xFF); HL=RGBColor(0xEB,0xFF,0x2C)
GREY=RGBColor(0x84,0x83,0x83); INK70=RGBColor(0x55,0x55,0x53); PALE=RGBColor(0xCF,0xCD,0xCA); FONT="Pretendard"
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height; BLANK=prs.slide_layouts[6]
def slide(bg):
    s=prs.slides.add_slide(BLANK); r=s.shapes.add_shape(1,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False; return s
def tb(s,l,t,w,h,a=MSO_ANCHOR.TOP):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); b.text_frame.word_wrap=True; b.text_frame.vertical_anchor=a; return b.text_frame
def para(tf,runs,size,bold=False,color=INK,space=6,first=False,lh=1.15):
    p=(tf.paragraphs[0] if first else tf.add_paragraph()); p.space_after=Pt(space); p.line_spacing=lh
    if isinstance(runs,str): runs=[(runs,color,bold)]
    for txt,col,bd in runs:
        r=p.add_run(); r.text=txt; f=r.font; f.name=FONT; f.size=Pt(size); f.bold=bd; f.color.rgb=col
    return p
def kicker(s,txt,light=True): para(tb(s,0.75,0.55,11.8,0.5),txt,16,True,(GREY if light else PALE),0,True)
def strap(s,pg,light=True):
    c=INK70 if light else PALE
    para(tb(s,0.75,7.0,11.8,0.4),f"WMBB · 엘리스랩 딥리서치 교육      2026.07.24        {pg}",11,False,c,0,True)

# 01 표지 (dark)
s=slide(INK); kicker(s,"리서치 보고서 · 딥리서치 에이전트 산출",False)
tf=tb(s,0.75,2.2,11.8,3,MSO_ANCHOR.MIDDLE)
para(tf,C["title"],38,True,WHITE,4,True,lh=1.05)
para(tb(s,0.78,5.2,11.8,0.6),C["meta"],17,False,PALE,0,True); strap(s,"01",False)
# 02 결론 (light)
s=slide(PAPER); kicker(s,"1 · 질문과 결론")
para(tb(s,0.75,1.1,11.8,1),"답부터 · 결론",30,True,INK,0,True)
box=s.shapes.add_shape(1,Inches(0.75),Inches(2.6),Inches(11.8),Inches(2.3))
box.fill.solid(); box.fill.fore_color.rgb=WHITE; box.line.color.rgb=RGBColor(0xE5,0xE5,0xE5); box.shadow.inherit=False
tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.4); tf.margin_top=Inches(0.3); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
para(tf,f"질문: {C['question']}",22,False,INK70,10,True)
para(tf,f"결론: {C['conclusion']}",25,True,INK,10)
para(tf,C["conclusion_sub"],17,False,INK70,0); strap(s,"02")
# 03 인사이트 (light)
s=slide(PAPER); kicker(s,"2 · 인사이트")
para(tb(s,0.75,1.1,11.8,1),"자료가 아니라 판단",30,True,INK,0,True)
y=2.35
for tag,head,body in C["insights"][:3]:
    tf=tb(s,0.75,y,11.8,1.3)
    para(tf,[(f"[{tag}] " if tag else "",INK,True),(head,INK,True)],19,True,INK,3,True,lh=1.15)
    para(tf,body,15,False,INK70,0); y+=1.45
strap(s,"03")
# 04 반대 근거 (dark)
s=slide(INK); kicker(s,"3 · 반대 근거",False)
tf=tb(s,0.75,2.2,11.8,3,MSO_ANCHOR.MIDDLE)
para(tf,[(C["counter_hl"],HL,True),(C["counter"],WHITE,True)],26,True,WHITE,14,True,lh=1.2)
para(tf,C["counter_sub"],19,False,PALE,0); strap(s,"04",False)
# 05 액션 (light)
s=slide(PAPER); kicker(s,"4 · 그래서 뭘 해야 하는가")
para(tb(s,0.75,1.1,11.8,1),"다음 액션",30,True,INK,0,True)
tf=tb(s,0.75,2.5,11.8,3.5)
for i,t in enumerate(C["actions"]):
    para(tf,[(f"0{i+1}   ",INK,True),(t,INK70,False)],19,False,INK70,16,i==0,lh=1.3)
strap(s,"05")
# 06 출처 (light)
s=slide(PAPER); kicker(s,"5 · "+C["sources_head"])
para(tb(s,0.75,1.1,11.8,1),"근거와 출처",30,True,INK,0,True)
tf=tb(s,0.75,2.4,11.8,4)
for i,t in enumerate(C["sources"]):
    para(tf,[(f"[{i+1}] ",INK,True),(t,INK70,False)],16,False,INK70,10,i==0,lh=1.3)
if C["dropped"]: para(tf,C["dropped"],14,False,GREY,0)
strap(s,"06")

prs.save(OUT)
print(f"{OUT} 생성 · {len(prs.slides._sldIdLst)}슬라이드 · (원본: {os.path.basename(md_path)})")
